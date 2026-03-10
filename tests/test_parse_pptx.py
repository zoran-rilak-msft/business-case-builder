"""Tests for scripts/parse_pptx.py."""

import json
import subprocess
from pathlib import Path

import pytest
from pptx import Presentation

PROJECT_ROOT = Path("/Users/zoranrilak/Documents/src/speckit/business-case-builder")
FIXTURE = PROJECT_ROOT / "tests" / "fixtures" / "sample-presentation.pptx"


def _run_parser(pptx_path: Path) -> subprocess.CompletedProcess:
    return subprocess.run(
        ["python", "scripts/parse_pptx.py", "--input", str(pptx_path)],
        capture_output=True,
        text=True,
        cwd=str(PROJECT_ROOT),
    )


class TestParsePptx:
    def test_extract_text_from_shapes(self):
        result = _run_parser(FIXTURE)
        assert result.returncode == 0, result.stderr
        data = json.loads(result.stdout)
        assert data["slide_count"] == 3
        assert data["slides"][0]["title"] == "SSO Integration Feature"

    def test_speaker_notes_extraction(self):
        result = _run_parser(FIXTURE)
        assert result.returncode == 0, result.stderr
        data = json.loads(result.stdout)
        slide2 = data["slides"][1]
        assert slide2["speaker_notes"], "Slide 2 should have speaker notes"
        assert "business drivers" in slide2["speaker_notes"]

    def test_table_data_extraction(self):
        result = _run_parser(FIXTURE)
        assert result.returncode == 0, result.stderr
        data = json.loads(result.stdout)
        slide3 = data["slides"][2]
        table = slide3["table_data"]
        assert table is not None, "Slide 3 should contain a table"
        assert len(table) >= 3
        assert len(table[0]) >= 3
        assert "Category" in table[0]

    def test_empty_presentation(self, tmp_path):
        pptx_path = tmp_path / "empty.pptx"
        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
        prs.save(str(pptx_path))

        result = _run_parser(pptx_path)
        assert result.returncode == 0, result.stderr
        data = json.loads(result.stdout)
        assert data["slide_count"] == 1

    def test_corrupt_file_handling(self, tmp_path):
        bad_path = tmp_path / "corrupt.pptx"
        bad_path.write_text("not a pptx")

        result = _run_parser(bad_path)
        assert result.returncode == 1
        assert "Error" in result.stderr or "error" in result.stderr.lower()

    def test_extracted_text_completeness(self):
        result = _run_parser(FIXTURE)
        assert result.returncode == 0, result.stderr
        data = json.loads(result.stdout)
        text = data["extracted_text"]
        for phrase in ("SSO Integration", "Key Benefits", "Impact Analysis"):
            assert phrase in text, f"extracted_text should contain '{phrase}'"
