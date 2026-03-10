"""Extract text content from PowerPoint presentations into structured JSON."""

import argparse
import json
import sys
from pathlib import Path

from pptx import Presentation

sys.path.insert(0, str(Path(__file__).resolve().parent.parent))
from scripts.models import PptxOutput, SlideContent  # noqa: E402


def _extract_slides(prs: Presentation) -> list[SlideContent]:
    slides: list[SlideContent] = []
    for idx, slide in enumerate(prs.slides, start=1):
        title = None
        body_parts: list[str] = []

        for shape in slide.shapes:
            if shape.has_table:
                continue
            if not shape.has_text_frame:
                continue
            if shape.is_placeholder and shape.placeholder_format.idx == 0:
                title = shape.text_frame.text
            else:
                text = shape.text_frame.text.strip()
                if text:
                    body_parts.append(text)

        table_data: list[list[str]] | None = None
        for shape in slide.shapes:
            if shape.has_table:
                table = shape.table
                table_data = [
                    [cell.text for cell in row.cells] for row in table.rows
                ]
                break

        speaker_notes = ""
        try:
            notes_slide = slide.notes_slide
            speaker_notes = notes_slide.notes_text_frame.text
        except Exception:
            pass

        body_text = "\n".join(body_parts)

        if not title and not body_text and not table_data and not speaker_notes:
            print(f"Warning: slide {idx} has no text content", file=sys.stderr)

        slides.append(
            SlideContent(
                slide_number=idx,
                title=title,
                body_text=body_text,
                speaker_notes=speaker_notes,
                table_data=table_data,
            )
        )
    return slides


def _build_extracted_text(slides: list[SlideContent]) -> str:
    parts: list[str] = []
    for s in slides:
        if s.title:
            parts.append(s.title)
        if s.body_text:
            parts.append(s.body_text)
        if s.speaker_notes:
            parts.append(s.speaker_notes)
        if s.table_data:
            for row in s.table_data:
                parts.append(" | ".join(row))
    return "\n".join(parts)


def main() -> None:
    parser = argparse.ArgumentParser(description="Extract text from a PPTX file.")
    parser.add_argument("--input", required=True, help="Path to the .pptx file")
    args = parser.parse_args()

    path = Path(args.input)
    if not path.exists():
        print(f"Error: file not found: {path}", file=sys.stderr)
        sys.exit(1)

    try:
        prs = Presentation(str(path))
    except Exception as exc:
        print(f"Error: could not open PPTX file: {exc}", file=sys.stderr)
        sys.exit(1)

    slides = _extract_slides(prs)
    extracted_text = _build_extracted_text(slides)

    output = PptxOutput(
        file_path=str(path),
        slide_count=len(slides),
        slides=slides,
        extracted_text=extracted_text,
    )

    total_chars = len(extracted_text)
    print(
        f"Parsed {output.slide_count} slides, extracted {total_chars} characters",
        file=sys.stderr,
    )
    print(json.dumps(output.model_dump(), indent=2))


if __name__ == "__main__":
    main()
