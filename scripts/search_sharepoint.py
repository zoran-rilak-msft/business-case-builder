"""Search SharePoint via Microsoft Graph and extract text from documents.

Two modes:
  --query "search terms" [--site <url>]   → search for documents, return metadata
  --download <drive-item-id>              → download document, extract full text
  --login                                 → authenticate and cache token
  --logout                                → clear cached token

Requires BCB_GRAPH_CLIENT_ID environment variable.
"""

from __future__ import annotations

import argparse
import io
import json
import sys
from pathlib import Path

import requests

sys.path.insert(0, str(Path(__file__).resolve().parent.parent))
from scripts.models import SharePointDocumentContent, SharePointSearchItem, SharePointSearchResult
from scripts.sharepoint_auth import clear_cache, get_graph_token

_GRAPH_BASE = "https://graph.microsoft.com/v1.0"

# Maximum number of search results to return
_MAX_RESULTS = 25


# ---------------------------------------------------------------------------
# Search mode
# ---------------------------------------------------------------------------


def _build_search_payload(query: str, site_scope: str | None) -> dict:
    """Build the JSON body for POST /search/query."""
    entity_request = {
        "entityTypes": ["driveItem"],
        "query": {"queryString": query},
        "from": 0,
        "size": _MAX_RESULTS,
    }
    if site_scope:
        # KQL site filter: path:{url} restricts results to that site
        entity_request["query"]["queryString"] = f"({query}) path:\"{site_scope}\""
    return {"requests": [entity_request]}


def _parse_search_response(raw: dict, query: str, site_scope: str | None) -> SharePointSearchResult:
    """Extract search results from Graph /search/query response."""
    items: list[SharePointSearchItem] = []

    for response_value in raw.get("value", []):
        for hit_container in response_value.get("hitsContainers", []):
            for hit in hit_container.get("hits", []):
                resource = hit.get("resource", {})
                name = resource.get("name", "Untitled")
                web_url = resource.get("webUrl", "")
                last_modified = resource.get("lastModifiedDateTime", "")
                size = resource.get("size", 0)

                # Extract site name from parent reference
                parent = resource.get("parentReference", {})
                site_name = parent.get("siteId", "").split(",")[0] if parent.get("siteId") else ""

                # Determine file type from extension
                ext = Path(name).suffix.lstrip(".").lower()
                file_type = ext if ext in ("docx", "pptx", "pdf", "xlsx") else ext

                # Extract drive item ID for download
                drive_item_id = resource.get("id", "")
                # Build a composite ID: driveId!itemId for download
                drive_id = parent.get("driveId", "")
                if drive_id and drive_item_id:
                    composite_id = f"{drive_id}!{drive_item_id}"
                else:
                    composite_id = drive_item_id

                # Extract snippet from summary
                snippet = hit.get("summary", "")

                items.append(SharePointSearchItem(
                    drive_item_id=composite_id,
                    title=name,
                    url=web_url,
                    site_name=site_name,
                    last_modified=last_modified,
                    file_type=file_type,
                    snippet=snippet,
                    size_bytes=size,
                ))

    return SharePointSearchResult(
        query=query,
        site_scope=site_scope,
        result_count=len(items),
        results=items,
    )


def search_sharepoint(query: str, site_scope: str | None = None) -> SharePointSearchResult:
    """Search SharePoint for documents matching the query."""
    token = get_graph_token()
    headers = {"Authorization": f"Bearer {token}", "Content-Type": "application/json"}
    payload = _build_search_payload(query, site_scope)

    resp = requests.post(f"{_GRAPH_BASE}/search/query", headers=headers, json=payload, timeout=30)

    if resp.status_code == 401:
        print(
            "Error: Authentication expired. "
            "Run: python scripts/search_sharepoint.py --login",
            file=sys.stderr,
        )
        sys.exit(1)
    if resp.status_code == 403:
        print(
            "Error: Insufficient permissions. "
            "Ensure the app has Sites.Read.All and Files.Read.All.",
            file=sys.stderr,
        )
        sys.exit(1)
    if resp.status_code == 429:
        retry_after = resp.headers.get("Retry-After", "60")
        print(
            f"Error: Rate limited by Graph API. Retry after {retry_after}s.",
            file=sys.stderr,
        )
        sys.exit(1)
    if not resp.ok:
        print(f"Error: Graph API returned {resp.status_code}: {resp.text[:500]}", file=sys.stderr)
        sys.exit(1)

    return _parse_search_response(resp.json(), query, site_scope)


# ---------------------------------------------------------------------------
# Download + extract mode
# ---------------------------------------------------------------------------


def _download_content(
    drive_item_id: str, token: str,
) -> tuple[bytes, str, str, str]:
    """Download file content via Graph API."""
    headers = {"Authorization": f"Bearer {token}"}

    # Parse composite ID: driveId!itemId
    if "!" in drive_item_id:
        drive_id, item_id = drive_item_id.split("!", 1)
        meta_url = f"{_GRAPH_BASE}/drives/{drive_id}/items/{item_id}"
    else:
        print("Error: Invalid drive_item_id format. Expected 'driveId!itemId'.", file=sys.stderr)
        sys.exit(1)

    # Get metadata first
    meta_resp = requests.get(meta_url, headers=headers, timeout=30)
    if not meta_resp.ok:
        print(f"Error: Could not fetch item metadata: {meta_resp.status_code}", file=sys.stderr)
        sys.exit(1)

    meta = meta_resp.json()
    filename = meta.get("name", "unknown")
    web_url = meta.get("webUrl", "")

    # Download content
    content_resp = requests.get(f"{meta_url}/content", headers=headers, timeout=120, stream=True)
    if not content_resp.ok:
        print(f"Error: Could not download file: {content_resp.status_code}", file=sys.stderr)
        sys.exit(1)

    return content_resp.content, filename, web_url, drive_item_id


def _extract_text_docx(content: bytes) -> tuple[str, int | None]:
    """Extract text from a Word document."""
    from docx import Document

    doc = Document(io.BytesIO(content))
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
    return "\n\n".join(paragraphs), None


def _extract_text_pptx(content: bytes) -> tuple[str, int | None]:
    """Extract text from a PowerPoint presentation."""
    from pptx import Presentation

    prs = Presentation(io.BytesIO(content))
    parts: list[str] = []
    for i, slide in enumerate(prs.slides, 1):
        slide_texts: list[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_texts.append(shape.text)
        if slide_texts:
            parts.append(f"[Slide {i}]\n" + "\n".join(slide_texts))
    return "\n\n".join(parts), len(prs.slides)


def _extract_text_pdf(content: bytes) -> tuple[str, int | None]:
    """Extract text from a PDF document."""
    from PyPDF2 import PdfReader

    reader = PdfReader(io.BytesIO(content))
    pages: list[str] = []
    for page in reader.pages:
        text = page.extract_text()
        if text and text.strip():
            pages.append(text.strip())
    return "\n\n".join(pages), len(reader.pages)


def _extract_text_xlsx(content: bytes) -> tuple[str, int | None]:
    """Extract text from an Excel workbook."""
    from openpyxl import load_workbook

    wb = load_workbook(io.BytesIO(content), read_only=True, data_only=True)
    parts: list[str] = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        rows: list[str] = []
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) if c is not None else "" for c in row]
            if any(cells):
                rows.append("\t".join(cells))
        if rows:
            parts.append(f"[Sheet: {sheet_name}]\n" + "\n".join(rows))
    wb.close()
    return "\n\n".join(parts), len(wb.sheetnames)


_EXTRACTORS = {
    "docx": ("python-docx", _extract_text_docx),
    "pptx": ("python-pptx", _extract_text_pptx),
    "pdf": ("PyPDF2", _extract_text_pdf),
    "xlsx": ("openpyxl", _extract_text_xlsx),
}


def download_and_extract(drive_item_id: str) -> SharePointDocumentContent:
    """Download a document from SharePoint and extract its text content."""
    token = get_graph_token()
    content, filename, web_url, item_id = _download_content(drive_item_id, token)

    ext = Path(filename).suffix.lstrip(".").lower()
    if ext not in _EXTRACTORS:
        print(
            f"Error: Unsupported file type '.{ext}'. Supported: .docx, .pptx, .pdf, .xlsx",
            file=sys.stderr,
        )
        sys.exit(1)

    method_name, extractor = _EXTRACTORS[ext]
    try:
        extracted_text, page_count = extractor(content)
    except Exception as exc:
        print(f"Error: Failed to extract text from {filename}: {exc}", file=sys.stderr)
        sys.exit(1)

    return SharePointDocumentContent(
        drive_item_id=item_id,
        title=filename,
        url=web_url,
        file_type=ext,
        extracted_text=extracted_text,
        page_count=page_count,
        extraction_method=method_name,
    )


# ---------------------------------------------------------------------------
# CLI
# ---------------------------------------------------------------------------


def _parse_args(argv: list[str] | None = None) -> argparse.Namespace:
    parser = argparse.ArgumentParser(
        description="Search SharePoint and extract document text via Microsoft Graph."
    )
    group = parser.add_mutually_exclusive_group(required=True)
    group.add_argument("--query", help="Search query (KQL syntax supported)")
    group.add_argument(
        "--download", metavar="DRIVE_ITEM_ID",
        help="Download and extract text from a document",
    )
    group.add_argument("--login", action="store_true", help="Authenticate and cache token")
    group.add_argument("--logout", action="store_true", help="Clear cached authentication token")

    parser.add_argument("--site", help="Scope search to a specific SharePoint site URL")
    return parser.parse_args(argv)


def main(argv: list[str] | None = None) -> None:
    args = _parse_args(argv)

    if args.login:
        token = get_graph_token()
        print(f"Authenticated successfully. Token length: {len(token)}", file=sys.stderr)
        return

    if args.logout:
        clear_cache()
        return

    if args.query:
        result = search_sharepoint(args.query, args.site)
        print(json.dumps(result.model_dump(), indent=2))
        print(f"Found {result.result_count} results for '{args.query}'", file=sys.stderr)
        return

    if args.download:
        result = download_and_extract(args.download)
        print(json.dumps(result.model_dump(), indent=2))
        print(f"Extracted text from {result.title} ({result.extraction_method})", file=sys.stderr)
        return


if __name__ == "__main__":
    main()
